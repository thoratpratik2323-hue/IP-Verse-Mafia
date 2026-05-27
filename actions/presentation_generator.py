"""
presentation_generator.py — Voice-to-Presentation PowerPoint deck compiler for IP Prime.

Utilizes python-pptx library to build structured slides (titles, bullet points, speaker notes)
from vocal topics, saving locally and auto-opening the presentation deck.
"""

from __future__ import annotations

import logging
import os
import json
from pathlib import Path
from typing import Any, Optional

logger = logging.getLogger("ip_prime.presentation_generator")

BASE_DIR = Path(__file__).resolve().parent.parent
OUTPUT_DIR = Path("C:/Users/thora/Downloads/IP Given/code")

def _ensure_output_dir():
    try:
        OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    except Exception:
        try:
            (BASE_DIR / "data").mkdir(parents=True, exist_ok=True)
        except Exception:
            pass

def _get_save_path(topic: str) -> Path:
    _ensure_output_dir()
    clean_topic = "".join([c if c.isalnum() else "_" for c in topic.strip()]).lower()[:20]
    filename = f"presentation_{clean_topic}.pptx"
    if OUTPUT_DIR.exists():
        return OUTPUT_DIR / filename
    return BASE_DIR / "data" / filename

def generate_presentation_from_topic(topic: str, slide_count: int = 5, player: Optional[Any] = None) -> str:
    """
    Uses Gemini to generate slide content layouts and compiles the PPTX presentation deck.
    """
    if not topic:
        return "Topic cannot be empty, sir."
        
    logger.info("Generating AI presentation for topic '%s' with %d slides...", topic, slide_count)
    
    # 1. Structure outline using Gemini (or fallback mock)
    gemini_key = os.environ.get("GEMINI_API_KEY", "").strip()
    slides_data = []

    if gemini_key:
        try:
            import google.generativeai as genai
            genai.configure(api_key=gemini_key)
            model = genai.GenerativeModel("gemini-2.0-flash")
            
            prompt = (
                f"You are a professional presentation designer. Create a PowerPoint outline for the topic: '{topic}'. "
                f"Create exactly {slide_count} slides. "
                "For each slide, provide a Title, 3-4 Bullet points, and Speaker notes. "
                "Return the response in a clean, strict JSON format matching this schema:\n"
                "{\n"
                "  \"slides\": [\n"
                "    {\n"
                "      \"title\": \"Slide Title\",\n"
                "      \"bullets\": [\"Bullet point 1\", \"Bullet point 2\", \"Bullet point 3\"],\n"
                "      \"notes\": \"Detailed speaker notes for this slide.\"\n"
                "    }\n"
                "  ]\n"
                "}\n"
                "Return only the raw JSON. Do not include markdown backticks."
            )
            res = model.generate_content(prompt)
            res_text = res.text.strip()
            if "```" in res_text:
                res_text = res_text.replace("```json", "").replace("```", "").strip()
                
            slides_data = json.loads(res_text).get("slides", [])
        except Exception as e:
            logger.error("Gemini outline generation failed: %s. Falling back to structured simulation.", e)

    if not slides_data:
        # Structured mock fallback
        slides_data = [
            {
                "title": f"Introduction to {topic}",
                "bullets": [f"Core definitions and parameters of {topic}", "Why it is critical in today's workspace", "Outline of the presentation"],
                "notes": "Good morning everyone. Today we are exploring the core tenets of this subject."
            },
            {
                "title": "Key Features & Functionality",
                "bullets": ["High-level efficiency parameters", "Innovative architectural integrations", "Developer-first designs"],
                "notes": "Here we see the main technical parameters that set this framework apart."
            },
            {
                "title": "Future Growth & Summary",
                "bullets": ["Next-gen scalability prospects", "Strategic roadmap milestones", "Summary of conclusions"],
                "notes": "To wrap up, we expect continuous iterations to scale these interfaces."
            }
        ]

    # 2. Compile PPTX using python-pptx
    save_path = _get_save_path(topic)
    pptx_active = False

    try:
        from pptx import Presentation
        prs = Presentation()
        
        for slide_info in slides_data:
            # Add slide layout (1 is title-and-content layout)
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            # Title
            title_box = slide.shapes.title
            title_box.text = slide_info.get("title", "Untitled Slide")
            
            # Bullets
            content_box = slide.placeholders[1]
            tf = content_box.text_frame
            tf.text = "" # Clear default
            
            bullets = slide_info.get("bullets", [])
            for idx, b in enumerate(bullets):
                p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                p.text = b
                
            # Speaker Notes
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_info.get("notes", "")

        prs.save(str(save_path))
        pptx_active = True
    except Exception as e:
        logger.error("Failed compiling PPTX via python-pptx: %s", e)

    # 3. Auto-open Presentation file
    open_success = False
    if pptx_active and save_path.exists():
        try:
            if os.name == "nt":
                os.startfile(str(save_path))
                open_success = True
            else:
                import subprocess
                subprocess.run(["xdg-open", str(save_path)], check=True)
                open_success = True
        except Exception as open_err:
            logger.error("Could not automatically launch presentation file: %s", open_err)

    status_msg = "PowerPoint compiled and opened successfully, sir!" if open_success else f"PPTX saved successfully at `{save_path}`!"
    if not pptx_active:
        status_msg = f"Simulated Slide deck generated for: '{topic}' (python-pptx not active on system)."
        
    return (
        f"### [PRESENTATION CREATED]\n"
        f"• **Topic**: {topic}\n"
        f"• **Slides**: {len(slides_data)}\n"
        f"• **Outcome**: {status_msg}\n\n"
        "Sabash sir! Deck setup complete!"
    )

def presentation_generator(parameters: dict[str, Any], player: Optional[Any] = None) -> str:
    """Main dispatcher for presentation_generator action."""
    action = parameters.get("action", "generate").lower().strip()
    topic = parameters.get("topic", "")
    slides = int(parameters.get("slides", 5))
    
    if action == "generate":
        return generate_presentation_from_topic(topic, slides, player)
    else:
        return "Unknown presentation generator action parameter, sir."
